"""Google Drive integration service."""

import io
from pathlib import Path
from typing import Any

from google.oauth2 import service_account
from google.oauth2.credentials import Credentials
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from loguru import logger

from app.config import settings


class DriveService:
    """Service for interacting with Google Drive API."""

    SCOPES = ["https://www.googleapis.com/auth/drive.readonly"]

    def __init__(self) -> None:
        """Initialize Google Drive service with auto-detected auth method."""
        credentials = self._get_credentials()
        self.service = build("drive", "v3", credentials=credentials)
    
    def _get_credentials(self):
        """Get credentials using configured auth method."""
        auth_method = settings.google_auth_method
        
        # Auto-detect: try OAuth first (for restricted orgs), then service account
        if auth_method == "auto":
            oauth_token = Path(settings.google_oauth_token_file)
            service_account_file = Path(settings.google_service_account_file)
            
            if oauth_token.exists():
                logger.info("Using OAuth user credentials")
                return self._get_oauth_credentials()
            elif service_account_file.exists():
                logger.info("Using service account credentials")
                return self._get_service_account_credentials()
            else:
                raise FileNotFoundError(
                    "No credentials found. Run 'python scripts/authenticate_drive.py' "
                    "or provide a service account file."
                )
        elif auth_method == "oauth":
            return self._get_oauth_credentials()
        else:  # service_account
            return self._get_service_account_credentials()
    
    def _get_service_account_credentials(self):
        """Get service account credentials."""
        return service_account.Credentials.from_service_account_file(
            settings.google_service_account_file,
            scopes=self.SCOPES,
        )
    
    def _get_oauth_credentials(self):
        """Get OAuth user credentials with auto-refresh."""
        token_path = Path(settings.google_oauth_token_file)
        
        if not token_path.exists():
            raise FileNotFoundError(
                f"OAuth token not found: {token_path}. "
                "Run 'python scripts/authenticate_drive.py' first."
            )
        
        credentials = Credentials.from_authorized_user_file(str(token_path), self.SCOPES)
        
        # Refresh if expired
        if credentials.expired and credentials.refresh_token:
            logger.info("Refreshing expired OAuth token...")
            credentials.refresh(Request())
            # Save refreshed token
            with open(token_path, "w") as f:
                f.write(credentials.to_json())
        
        return credentials

    def list_files(self, folder_id: str | None = None) -> list[dict[str, Any]]:
        """
        List all files in a folder (non-recursive).
        
        Args:
            folder_id: Google Drive folder ID (uses config default if not provided)
            
        Returns:
            List of file metadata dictionaries
        """
        folder_id = folder_id or settings.google_drive_folder_id
        
        query = f"'{folder_id}' in parents and trashed = false"
        
        all_files = []
        page_token = None
        
        while True:
            results = self.service.files().list(
                q=query,
                pageSize=100,
                pageToken=page_token,
                fields="nextPageToken, files(id, name, mimeType, webViewLink, modifiedTime)",
            ).execute()
            
            all_files.extend(results.get("files", []))
            page_token = results.get("nextPageToken")
            
            if not page_token:
                break
        
        logger.info(f"Found {len(all_files)} files in folder {folder_id}")
        return all_files

    def list_files_recursive(self, folder_id: str | None = None, path: str = "") -> list[dict[str, Any]]:
        """
        Recursively list all files in a folder and its subfolders.
        
        Args:
            folder_id: Google Drive folder ID (uses config default if not provided)
            path: Current path for logging (internal use)
            
        Returns:
            List of file metadata dictionaries with 'path' field added
        """
        folder_id = folder_id or settings.google_drive_folder_id
        
        all_files = []
        items = self.list_files(folder_id)
        
        for item in items:
            item_path = f"{path}/{item['name']}" if path else item['name']
            
            if item.get("mimeType") == "application/vnd.google-apps.folder":
                # Recursively process subfolder
                logger.info(f"Entering folder: {item_path}")
                subfolder_files = self.list_files_recursive(item["id"], item_path)
                all_files.extend(subfolder_files)
            else:
                # Add path to file info
                item["path"] = item_path
                all_files.append(item)
        
        return all_files

    def get_file_content(self, file_id: str) -> str:
        """
        Download and return file content as text.
        
        Supports Google Docs, Sheets, plain text, Word docs, and PowerPoint files.
        
        Args:
            file_id: Google Drive file ID
            
        Returns:
            File content as string
        """
        # Get file metadata
        file_meta = self.service.files().get(
            fileId=file_id,
            fields="mimeType, name",
        ).execute()
        
        mime_type = file_meta.get("mimeType", "")
        file_name = file_meta.get("name", "")
        
        # Skip folders
        if mime_type == "application/vnd.google-apps.folder":
            logger.debug(f"Skipping folder: {file_name}")
            return ""
        
        # Handle Google Docs - export as plain text
        if mime_type == "application/vnd.google-apps.document":
            request = self.service.files().export_media(
                fileId=file_id,
                mimeType="text/plain",
            )
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            content = buffer.getvalue().decode("utf-8")
            
        # Handle Google Sheets - export as CSV
        elif mime_type == "application/vnd.google-apps.spreadsheet":
            request = self.service.files().export_media(
                fileId=file_id,
                mimeType="text/csv",
            )
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            content = buffer.getvalue().decode("utf-8")
            
        # Handle Google Slides - export as plain text
        elif mime_type == "application/vnd.google-apps.presentation":
            request = self.service.files().export_media(
                fileId=file_id,
                mimeType="text/plain",
            )
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            content = buffer.getvalue().decode("utf-8")
            
        # Handle Word documents (.docx)
        elif mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_name.endswith(".docx"):
            request = self.service.files().get_media(fileId=file_id)
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buffer.seek(0)
            content = self._extract_docx_text(buffer)
            
        # Handle PowerPoint (.pptx)
        elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or file_name.endswith(".pptx"):
            request = self.service.files().get_media(fileId=file_id)
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buffer.seek(0)
            content = self._extract_pptx_text(buffer)
            
        # Handle PDF files
        elif mime_type == "application/pdf" or file_name.endswith(".pdf"):
            request = self.service.files().get_media(fileId=file_id)
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            buffer.seek(0)
            content = self._extract_pdf_text(buffer)
            
        # Handle plain text files
        else:
            request = self.service.files().get_media(fileId=file_id)
            buffer = io.BytesIO()
            downloader = MediaIoBaseDownload(buffer, request)
            done = False
            while not done:
                _, done = downloader.next_chunk()
            try:
                content = buffer.getvalue().decode("utf-8")
            except UnicodeDecodeError:
                content = buffer.getvalue().decode("latin-1", errors="ignore")
        
        logger.debug(f"Downloaded file: {file_name} ({len(content)} chars)")
        return content
    
    def _extract_docx_text(self, file_buffer: io.BytesIO) -> str:
        """Extract text from a Word document."""
        try:
            from docx import Document
            doc = Document(file_buffer)
            paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
            return "\n\n".join(paragraphs)
        except Exception as e:
            logger.warning(f"Failed to extract docx text: {e}")
            return ""
    
    def _extract_pptx_text(self, file_buffer: io.BytesIO) -> str:
        """Extract text from a PowerPoint presentation."""
        try:
            from pptx import Presentation
            prs = Presentation(file_buffer)
            text_parts = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                if slide_text:
                    text_parts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.warning(f"Failed to extract pptx text: {e}")
            return ""
    
    def _extract_pdf_text(self, file_buffer: io.BytesIO) -> str:
        """Extract text from a PDF file."""
        try:
            import pypdf
            reader = pypdf.PdfReader(file_buffer)
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text.strip():
                    text_parts.append(text)
            return "\n\n".join(text_parts)
        except Exception as e:
            logger.warning(f"Failed to extract PDF text: {e}")
            return ""

    def get_file_url(self, file_id: str) -> str:
        """Get the web view URL for a file."""
        return f"https://drive.google.com/file/d/{file_id}/view"
